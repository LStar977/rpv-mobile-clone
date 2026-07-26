"""Render a .pptx to HTML for visual QA.

LibreOffice is broken in this sandbox, so instead of converting the deck we read
the generated file back with python-pptx and draw exactly what it contains —
every shape's real position, size, fill and text. Fonts are mapped to their
metric-compatible clones (Calibri->Carlito, Cambria->Caladea) so text wraps at
the same width PowerPoint would use, which is what makes overflow checks
trustworthy.
"""
import base64
import html
import sys

from pptx import Presentation
from pptx.util import Emu

EMU_IN = 914400.0
PX_IN = 96.0

FONT_MAP = {
    'Calibri': 'Carlito, Calibri, sans-serif',
    'Cambria': 'Caladea, Cambria, serif',
    'Courier New': '"Courier New", monospace',
}


def px(emu):
    return (emu or 0) / EMU_IN * PX_IN


def color_of(fmt):
    """Best-effort solid colour as #rrggbb, or None."""
    try:
        if fmt.type is None:
            return None
        c = fmt.fore_color
        if c.type is not None and str(c.type).startswith('MSO_THEME'):
            return None
        return '#' + str(c.rgb)
    except Exception:
        return None


def line_of(shape):
    try:
        ln = shape.line
        if ln.fill.type is None or str(ln.fill.type) == 'MSO_FILL_TYPE.BACKGROUND (5)':
            return None
        w = px(ln.width) if ln.width else 1
        return '%.2fpx solid #%s' % (max(w, 0.75), str(ln.color.rgb))
    except Exception:
        return None


def run_html(run):
    f = run.font
    styles = []
    if f.size:
        styles.append('font-size:%.2fpx' % (f.size.pt * PX_IN / 72))
    if f.bold:
        styles.append('font-weight:700')
    if f.italic:
        styles.append('font-style:italic')
    try:
        if f.color and f.color.type is not None:
            styles.append('color:#' + str(f.color.rgb))
    except Exception:
        pass
    if f.name:
        styles.append('font-family:%s' % FONT_MAP.get(f.name, f.name))
    return '<span style="%s">%s</span>' % (';'.join(styles), html.escape(run.text))


ALIGN = {2: 'center', 3: 'right', 4: 'justify'}  # PP_ALIGN: LEFT=1
VANCHOR = {3: 'center', 4: 'flex-end'}  # MSO_ANCHOR: TOP=1, MIDDLE=3, BOTTOM=4


def shape_html(sh):
    out = []
    left, top = px(sh.left), px(sh.top)
    w, h = px(sh.width), px(sh.height)
    base = 'position:absolute;left:%.2fpx;top:%.2fpx;width:%.2fpx;height:%.2fpx;' % (left, top, w, h)

    # picture
    if sh.shape_type == 13 or getattr(sh, 'image', None) is not None:
        try:
            b64 = base64.b64encode(sh.image.blob).decode()
            out.append('<img style="%s" src="data:image/png;base64,%s"/>' % (base, b64))
            return out
        except Exception:
            pass

    if getattr(sh, 'has_chart', False):
        out.append(
            '<div style="%sborder:1px dashed #999;display:flex;align-items:center;'
            'justify-content:center;color:#999;font:12px sans-serif">[native chart]</div>' % base
        )
        return out

    # shape body
    box = base
    fill = color_of(sh.fill) if hasattr(sh, 'fill') else None
    if fill:
        box += 'background:%s;' % fill
    ln = line_of(sh)
    if ln:
        box += 'border:%s;' % ln
    prst = ''
    try:
        prst = sh._element.spPr.find(
            '{http://schemas.openxmlformats.org/drawingml/2006/main}prstGeom'
        ).get('prst') or ''
    except Exception:
        pass
    if prst == 'ellipse':
        box += 'border-radius:50%;'
    elif 'round' in prst.lower():
        box += 'border-radius:6px;'

    inner = ''
    if sh.has_text_frame and sh.text_frame.text.strip():
        tf = sh.text_frame
        mt = px(tf.margin_top) if tf.margin_top is not None else 0
        mb = px(tf.margin_bottom) if tf.margin_bottom is not None else 0
        ml = px(tf.margin_left) if tf.margin_left is not None else 0
        mr = px(tf.margin_right) if tf.margin_right is not None else 0
        va = VANCHOR.get(
            tf.vertical_anchor if tf.vertical_anchor is None else int(tf.vertical_anchor), 'flex-start'
        )
        paras = []
        for p in tf.paragraphs:
            if not p.runs:
                paras.append('<p style="margin:0;height:0.5em"></p>')
                continue
            al = ALIGN.get(p.alignment if p.alignment is None else int(p.alignment), 'left')
            bullet = ''
            if p._pPr is not None and p._pPr.find(
                '{http://schemas.openxmlformats.org/drawingml/2006/main}buChar'
            ) is not None:
                bullet = '• '
            sa = 0
            try:
                if p.space_after:
                    sa = p.space_after.pt * PX_IN / 72
            except Exception:
                pass
            paras.append(
                '<p style="margin:0 0 %.1fpx 0;text-align:%s;line-height:1.22">%s%s</p>'
                % (sa, al, bullet, ''.join(run_html(r) for r in p.runs))
            )
        inner = (
            '<div style="position:absolute;left:%.2fpx;top:%.2fpx;width:%.2fpx;height:%.2fpx;'
            'display:flex;flex-direction:column;justify-content:%s;overflow:visible">%s</div>'
            % (left + ml, top + mt, max(w - ml - mr, 1), max(h - mt - mb, 1), va, ''.join(paras))
        )

    out.append('<div style="%s"></div>' % box)
    if inner:
        out.append(inner)
    return out


def main(path, out_html):
    prs = Presentation(path)
    sw, sh_ = px(prs.slide_width), px(prs.slide_height)

    body = []
    for i, slide in enumerate(prs.slides, 1):
        bg = '#FFFFFF'
        try:
            c = slide.background.fill.fore_color.rgb
            bg = '#' + str(c)
        except Exception:
            pass
        parts = []
        for shp in slide.shapes:
            parts.extend(shape_html(shp))
        body.append(
            '<div class="wrap"><div class="num">Slide %d</div>'
            '<div class="slide" id="s%d" style="width:%.1fpx;height:%.1fpx;background:%s">%s</div></div>'
            % (i, i, sw, sh_, bg, ''.join(parts))
        )

    doc = (
        '<html><head><meta charset="utf-8"><style>'
        'body{margin:0;background:#555;font-family:Carlito,sans-serif}'
        '.wrap{margin:18px auto;width:%.1fpx}'
        '.num{color:#fff;font:12px sans-serif;margin-bottom:4px}'
        '.slide{position:relative;overflow:hidden}'
        '</style></head><body>%s</body></html>' % (sw, ''.join(body))
    )
    with open(out_html, 'w') as f:
        f.write(doc)
    print('wrote %s (%d slides)' % (out_html, len(body)))


if __name__ == '__main__':
    main(sys.argv[1], sys.argv[2])
